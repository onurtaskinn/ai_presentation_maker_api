from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR

# Create a new presentation
prs = Presentation()

# Add a slide with title and content layout
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)

# Set title and content
title = slide.shapes.title
title.text = "Your Title Here"

content = slide.placeholders[1]
content.text = "Your content here"

# Add charts, tables, images
# Example: Add an image
slide.shapes.add_picture('image.png', Inches(5), Inches(5), Inches(5), Inches(5))

# Save the presentation
prs.save('presentation.pptx')