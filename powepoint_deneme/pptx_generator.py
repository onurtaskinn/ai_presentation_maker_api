import warnings
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# Suppress zipfile warnings
warnings.filterwarnings('ignore', category=UserWarning, module='zipfile')

class PowerPointGenerator:
    """
    PowerPoint generator adapted for the AI presentation system
    """
    
    def __init__(self):
        self.prs = None
    
    def create_presentation_safely(self):
        """Create presentation without corruption issues"""
        try:
            # Always start fresh to avoid template issues
            self.prs = Presentation()
            
            # Set slide size AFTER creating presentation (16:9 format)
            self.prs.slide_width = Inches(16)
            self.prs.slide_height = Inches(9)
            
            return True
        except Exception as e:
            print(f"Error creating presentation: {e}")
            return False
    
    def get_local_image_path(self, presentation_id, slide_number):
        """Get local image path that was already downloaded"""
        try:
            # Use the same path structure as your download_image_to_local function
            image_path = f"images/{presentation_id}/slide_{slide_number}.jpg"
            
            if os.path.exists(image_path):
                return image_path
            else:
                print(f"Local image not found: {image_path}")
                return None
        except Exception as e:
            print(f"Error getting local image path: {e}")
            return None
    
    def add_title_slide(self, presentation_title):
        """Add title slide with presentation title"""
        try:
            if len(self.prs.slide_layouts) == 0:
                print("No slide layouts available")
                return False
                
            title_slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
            
            # Add title
            if title_slide.shapes.title:
                title_slide.shapes.title.text = presentation_title
                
                # Style the title
                title_shape = title_slide.shapes.title
                title_shape.text_frame.word_wrap = True
                
                # Format title text
                for paragraph in title_shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(44)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 32, 96)  # Dark blue
            
            # Add subtitle if placeholder exists
            if len(title_slide.placeholders) > 1:
                subtitle = title_slide.placeholders[1]
                subtitle.text = 'AI Generated Presentation'
                
                # Style subtitle
                for paragraph in subtitle.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(24)
                        run.font.color.rgb = RGBColor(89, 89, 89)  # Gray
            
            print("✓ Title slide added successfully")
            return True
            
        except Exception as e:
            print(f"Error adding title slide: {e}")
            return False
    
    def add_content_slide(self, slide_data, presentation_id):
        """Add content slide with title, text, and image"""
        try:
            # Use layout 1 (Title and Content) if available, otherwise layout 0
            layout_index = 1 if len(self.prs.slide_layouts) > 1 else 0
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_index])
            
            slide_number = slide_data["number"]
            slide_title = slide_data["title"]
            onscreen_text = slide_data["content"].slide_onscreen_text.text_list
            
            # Add slide title
            if slide.shapes.title:
                slide.shapes.title.text = slide_title
                
                # Style the title
                title_shape = slide.shapes.title
                for paragraph in title_shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(32)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 32, 96)
            
            # Add content text
            self.add_text_content(slide, onscreen_text)
            
            # Add image (using local file)
            self.add_image_to_slide(slide, presentation_id, slide_number)
            
            print(f"✓ Content slide {slide_number} added successfully")
            return True
            
        except Exception as e:
            print(f"Error adding content slide: {e}")
            return False
    
    def add_text_content(self, slide, onscreen_text_list):
        """Add text content to slide"""
        try:
            # Find content placeholder or create text box
            content_placeholder = None
            
            # Look for content placeholder
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 1:  # Content placeholder
                    content_placeholder = placeholder
                    break
            
            if content_placeholder is None:
                # Create text box if no content placeholder
                left = Inches(1)
                top = Inches(2)
                width = Inches(7)
                height = Inches(5)
                content_placeholder = slide.shapes.add_textbox(left, top, width, height)
            
            # Clear existing content
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            # Add text from list
            if onscreen_text_list:
                for i, text_item in enumerate(onscreen_text_list):
                    if i == 0:
                        # First paragraph
                        p = text_frame.paragraphs[0]
                    else:
                        # Additional paragraphs
                        p = text_frame.add_paragraph()
                    
                    p.text = text_item
                    p.level = 0
                    
                    # Style the text
                    for run in p.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(18)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                    
                    # Add some spacing between bullet points
                    p.space_after = Pt(6)
            
            return True
            
        except Exception as e:
            print(f"Error adding text content: {e}")
            return False
    
    def add_image_to_slide(self, slide, presentation_id, slide_number):
        """Add image to slide from local file"""
        try:
            # Get local image path
            local_image_path = self.get_local_image_path(presentation_id, slide_number)
            
            if local_image_path and os.path.exists(local_image_path):
                # Position image on the right side of the slide
                left = Inches(9)
                top = Inches(2)
                max_width = Inches(6)
                max_height = Inches(5)
                
                # Add image with size constraints
                pic = slide.shapes.add_picture(
                    local_image_path, 
                    left, 
                    top, 
                    width=max_width
                )
                
                # Adjust height if needed
                if pic.height > max_height:
                    pic.height = max_height
                
                print(f"✓ Image added to slide {slide_number}")
                return True
            else:
                print(f"⚠ Could not find local image for slide {slide_number}")
                return False
                
        except Exception as e:
            print(f"Error adding image to slide: {e}")
            return False
    
    def cleanup_temp_images(self):
        """Clean up temporary downloaded images - No longer needed since images are pre-downloaded"""
        # Images are already handled by the main system's cleanup
        # This method is kept for compatibility but does nothing
        pass
    
    def save_presentation_safely(self, filename, presentation_id):
        """Save presentation with validation"""
        try:
            # Ensure output directory exists
            output_dir = "_outputs"
            os.makedirs(output_dir, exist_ok=True)
            
            # Ensure filename has .pptx extension
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            # Full path
            full_path = os.path.join(output_dir, filename)
            
            # Save presentation
            self.prs.save(full_path)
            print(f"✓ Presentation saved successfully: {full_path}")
            
            # Validate saved file
            if os.path.exists(full_path):
                file_size = os.path.getsize(full_path)
                print(f"✓ File size: {file_size:,} bytes")
                
                # Basic validation - PPTX files should be > 10KB
                if file_size < 10000:
                    print("⚠ Warning: File seems unusually small")
                    return False
                else:
                    print("✓ File size looks good")
            
            return full_path
            
        except Exception as e:
            print(f"Error saving presentation: {e}")
            return None
    
    def validate_presentation_structure(self):
        """Validate presentation structure before saving"""
        try:
            print("\n=== Presentation Validation ===")
            
            if not self.prs:
                print("❌ No presentation created")
                return False
            
            # Check basic structure
            print(f"✓ Slide count: {len(self.prs.slides)}")
            print(f"✓ Layout count: {len(self.prs.slide_layouts)}")
            print(f"✓ Slide size: {self.prs.slide_width.inches}x{self.prs.slide_height.inches} inches")
            
            # Check each slide
            for i, slide in enumerate(self.prs.slides):
                print(f"✓ Slide {i+1}: {len(slide.shapes)} shapes")
                
                # Check for title
                if slide.shapes.title:
                    title_text = slide.shapes.title.text if slide.shapes.title.text else "(empty)"
                    print(f"  - Title: {title_text[:50]}{'...' if len(title_text) > 50 else ''}")
            
            print("✓ Validation completed successfully")
            return True
            
        except Exception as e:
            print(f"❌ Validation failed: {e}")
            return False


def create_presentation_from_data(presentation_data):
    """
    Main function to create PowerPoint presentation from AI-generated data
    
    Args:
        presentation_data: Dictionary containing presentation data with structure:
        {
            "id": str,
            "title": str,
            "slide_count": int,
            "slides": [
                {
                    "number": int,
                    "title": str,
                    "content": {
                        "slide_onscreen_text": {"text_list": [str, ...]},
                        ...
                    },
                    "image_url": str
                },
                ...
            ]
        }
    
    Returns:
        str: Path to created PowerPoint file, or None if failed
    """
    
    generator = PowerPointGenerator()
    
    try:
        print("=== PowerPoint Generation Started ===")
        
        # Extract data
        presentation_id = presentation_data["id"]
        presentation_title = presentation_data["title"]
        slide_count = presentation_data["slide_count"]
        slides = presentation_data["slides"]
        
        print(f"Creating presentation: {presentation_title}")
        print(f"Slides to generate: {slide_count}")
        
        # Create presentation
        if not generator.create_presentation_safely():
            return None
        
        # Add title slide
        if not generator.add_title_slide(presentation_title):
            return None
        
        # Add content slides
        for slide_data in slides:
            if not generator.add_content_slide(slide_data, presentation_id):
                print(f"⚠ Warning: Failed to add slide {slide_data.get('number', '?')}")
                # Continue with other slides
        
        # Validate structure
        if not generator.validate_presentation_structure():
            print("❌ Presentation structure validation failed")
            return None
        
        # Generate filename
        safe_title = "".join(c for c in presentation_title if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_title = safe_title.replace(' ', '_')
        filename = f"{presentation_id}_{safe_title}.pptx"
        
        # Save presentation
        file_path = generator.save_presentation_safely(filename, presentation_id)
        
        if file_path:
            print(f"\n🎉 SUCCESS: PowerPoint presentation created!")
            print(f"📁 File: {file_path}")
            return file_path
        else:
            print("\n❌ ERROR: Failed to save presentation")
            return None
            
    except Exception as e:
        print(f"\n❌ CRITICAL ERROR: {e}")
        return None
        
    finally:
        # Cleanup temporary files
        generator.cleanup_temp_images()

