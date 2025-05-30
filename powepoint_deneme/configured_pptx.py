from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import pandas as pd

# ==========================================
# 1. PRESENTATION INITIALIZATION CHOICES
# ==========================================

# Choice 1: Start with blank presentation or template
prs = Presentation()  # Blank presentation
# prs = Presentation('template.pptx')  # Use existing template

# Choice 2: Set slide size (optional)
# Default is usually fine, but you can customize
prs.slide_width = Inches(16)  # Wide format
prs.slide_height = Inches(9)

# ==========================================
# 2. SLIDE LAYOUT CHOICES
# ==========================================

# Available built-in layouts (indexes may vary by template):
layout_choices = {
    0: "Title Slide",
    1: "Title and Content", 
    2: "Section Header",
    3: "Two Content",
    4: "Comparison",
    5: "Title Only",
    6: "Blank",
    7: "Content with Caption",
    8: "Picture with Caption"
}

# Create slides with different layouts
title_slide_layout = prs.slide_layouts[0]  # Title slide
content_layout = prs.slide_layouts[1]      # Title and content
blank_layout = prs.slide_layouts[6]        # Blank slide

# ==========================================
# 3. TEXT AND TYPOGRAPHY CHOICES
# ==========================================

def add_text_slide_with_choices(prs):
    slide = prs.slides.add_slide(content_layout)
    
    # Title choices
    title = slide.shapes.title
    title.text = "Text Formatting Choices"
    
    # Content text box
    content = slide.placeholders[1]
    tf = content.text_frame
    
    # Text alignment choices
    p = tf.paragraphs[0]
    p.text = "Left aligned text"
    p.alignment = PP_ALIGN.LEFT
    
    # Add more paragraphs with different alignments
    p2 = tf.add_paragraph()
    p2.text = "Center aligned text"
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf.add_paragraph()
    p3.text = "Right aligned text"
    p3.alignment = PP_ALIGN.RIGHT
    
    # Font choices for individual runs
    run = p.runs[0]
    run.font.name = 'Arial'  # Font family choices: Arial, Calibri, Times New Roman, etc.
    run.font.size = Pt(18)   # Font size in points
    run.font.bold = True     # Bold formatting
    run.font.italic = False  # Italic formatting
    run.font.color.rgb = RGBColor(255, 0, 0)  # RGB color
    
    return slide

# ==========================================
# 7. IMAGE AND MEDIA CHOICES
# ==========================================

def add_image_slide(prs, image_path):
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Image Integration"
    
    # Image positioning choices
    left = Inches(1)
    top = Inches(2)
    
    # Size choices
    # Option 1: Specify width, height auto-calculated
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(6))
    
    # Option 2: Specify both dimensions
    # pic = slide.shapes.add_picture(image_path, left, top, Inches(6), Inches(4))
    
    return slide

# ==========================================
# 8. SLIDE TRANSITION AND ANIMATION CHOICES
# ==========================================

def set_slide_transitions(slide):
    # Note: python-pptx has limited animation support
    # Most animations need to be set in PowerPoint directly
    pass

# ==========================================
# 9. PRESENTATION ASSEMBLY AND EXPORT
# ==========================================

def create_presentation_with_data(data_dict):
    """
    Main function to create presentation with various choices
    
    data_dict should contain:
    - 'title': str
    - 'dataframe': pandas DataFrame
    - 'image_path': str (optional)
    """
    
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = data_dict.get('title', 'Data Presentation')
    subtitle = title_slide.placeholders[1]
    subtitle.text = 'Generated with Python-PPTX'
    
    # Add content slides
    add_text_slide_with_choices(prs)
    
    if 'image_path' in data_dict:
        add_image_slide(prs, data_dict['image_path'])
    
    return prs

# ==========================================
# 10. SAVE OPTIONS
# ==========================================

def save_presentation(prs, filename):
    """Save with different format choices"""
    
    # Standard save
    prs.save(filename)
    
    # You can also save with specific path
    # prs.save(f'/path/to/presentations/{filename}')

# ==========================================
# EXAMPLE USAGE
# ==========================================

if __name__ == "__main__":
    # Sample data
    sample_data = pd.DataFrame({
        'Category': ['A', 'B', 'C', 'D'],
        'Values': [23, 45, 56, 78],
        'Description': ['First', 'Second', 'Third', 'Fourth']
    })
    
    # Create presentation
    data_dict = {
        'title': 'My Data Presentation',
        'dataframe': sample_data,
        # 'image_path': 'path/to/image.png'  # Optional
    }
    
    presentation = create_presentation_with_data(data_dict)
    save_presentation(presentation, 'my_presentation.pptx')
    print("Presentation created successfully!")

# ==========================================
# KEY CHOICES SUMMARY
# ==========================================

"""
MAJOR CHOICE CATEGORIES:

1. PRESENTATION STRUCTURE:
   - Blank vs template-based
   - Slide layouts (Title, Content, Blank, etc.)
   - Slide dimensions

2. TEXT FORMATTING:
   - Font families, sizes, colors
   - Alignment (left, center, right, justify)
   - Bold, italic, underline
   - Text positioning and text boxes

3. VISUAL ELEMENTS:
   - Shapes (rectangles, circles, arrows, etc.)
   - Colors (RGB, theme colors)
   - Borders and fills
   - Transparency and effects

4. DATA PRESENTATION:
   - Tables with custom formatting
   - Charts (column, bar, line, pie, scatter)
   - Data arrangement and styling

5. MEDIA INTEGRATION:
   - Images with positioning and sizing
   - (Limited video/audio support)

6. LAYOUT AND POSITIONING:
   - Precise positioning with Inches/Pt
   - Relative positioning
   - Shape grouping and layering

7. EXPORT OPTIONS:
   - File naming and paths
   - Compression settings (limited)
"""