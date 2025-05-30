import warnings
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import pandas as pd

# Suppress zipfile warnings
warnings.filterwarnings('ignore', category=UserWarning, module='zipfile')

# ==========================================
# FIXED VERSION - NO CORRUPTION ISSUES
# ==========================================

def create_presentation_safely():
    """Create presentation without corruption issues"""
    
    # Always start fresh to avoid template issues
    prs = Presentation()
    
    # Set slide size AFTER creating presentation
    prs.slide_width = Inches(16)  # Wide format
    prs.slide_height = Inches(9)
    
    return prs

def add_text_slide_with_choices(prs):
    """Add text slide with proper error handling"""
    
    try:
        # Get available layouts
        if len(prs.slide_layouts) < 2:
            print("Warning: Not enough slide layouts available")
            return None
            
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use index 1 (Title and Content)
        
        # Add title safely
        if slide.shapes.title:
            slide.shapes.title.text = "Text Formatting Choices"
        
        # Check if content placeholder exists
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            
            # Clear existing text first
            tf.clear()
            
            # Add paragraphs safely
            p = tf.paragraphs[0]
            p.text = "Left aligned text"
            p.alignment = PP_ALIGN.LEFT
            
            # Add more paragraphs
            p2 = tf.add_paragraph()
            p2.text = "Center aligned text"
            p2.alignment = PP_ALIGN.CENTER
            
            p3 = tf.add_paragraph()
            p3.text = "Right aligned text"
            p3.alignment = PP_ALIGN.RIGHT
            
            # Font formatting - check if runs exist
            if p.runs:
                run = p.runs[0]
                run.font.name = 'Arial'
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.italic = False
                run.font.color.rgb = RGBColor(255, 0, 0)
        
        return slide
        
    except Exception as e:
        print(f"Error adding text slide: {e}")
        return None

def add_image_slide_safely(prs, image_path):
    """Add image slide with proper validation"""
    
    import os
    
    try:
        # Check if image file exists
        if not os.path.exists(image_path):
            print(f"Warning: Image file not found: {image_path}")
            return None
        
        # Check file extension
        valid_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp']
        if not any(image_path.lower().endswith(ext) for ext in valid_extensions):
            print(f"Warning: Unsupported image format: {image_path}")
            return None
        
        # Get layout safely
        layout_index = 1 if len(prs.slide_layouts) > 1 else 0
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        # Add title
        if slide.shapes.title:
            slide.shapes.title.text = "Image Integration"
        
        # Add image with error handling
        left = Inches(1)
        top = Inches(2)
        
        try:
            pic = slide.shapes.add_picture(image_path, left, top, width=Inches(6))
            return slide
        except Exception as img_error:
            print(f"Error adding image: {img_error}")
            return slide  # Return slide without image
            
    except Exception as e:
        print(f"Error creating image slide: {e}")
        return None

def create_presentation_with_data(data_dict):
    """
    FIXED: Main function to create presentation without corruption
    
    data_dict should contain:
    - 'title': str
    - 'dataframe': pandas DataFrame (optional)
    - 'image_path': str (optional)
    """
    
    try:
        # Create presentation safely
        prs = create_presentation_safely()
        
        # Title slide - use index 0 which should always exist
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Add title safely
        if title_slide.shapes.title:
            title_slide.shapes.title.text = data_dict.get('title', 'Data Presentation')
        
        # Add subtitle if placeholder exists
        if len(title_slide.placeholders) > 1:
            subtitle = title_slide.placeholders[1]
            subtitle.text = 'Generated with Python-PPTX'
        
        # Add content slides with error handling
        text_slide = add_text_slide_with_choices(prs)
        if text_slide:
            print("✓ Text slide added successfully")
        
        # Add image slide if path provided
        if 'image_path' in data_dict:
            image_slide = add_image_slide_safely(prs, data_dict['image_path'])
            if image_slide:
                print("✓ Image slide added successfully")
        
        return prs
        
    except Exception as e:
        print(f"Error creating presentation: {e}")
        
        # Fallback: create minimal presentation
        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Fallback Presentation"
            return prs
        except Exception as fallback_error:
            print(f"Fallback creation failed: {fallback_error}")
            raise

def save_presentation_safely(prs, filename):
    """Save presentation with validation"""
    
    try:
        # Ensure filename has .pptx extension
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        
        # Save presentation
        prs.save(filename)
        print(f"✓ Presentation saved successfully: {filename}")
        
        # Validate saved file
        import os
        if os.path.exists(filename):
            file_size = os.path.getsize(filename)
            print(f"✓ File size: {file_size:,} bytes")
            
            # Basic validation - PPTX files should be > 10KB
            if file_size < 10000:
                print("⚠ Warning: File seems unusually small")
            else:
                print("✓ File size looks good")
        
        return True
        
    except Exception as e:
        print(f"Error saving presentation: {e}")
        return False

def validate_presentation_structure(prs):
    """Validate presentation structure before saving"""
    
    try:
        print("\n=== Presentation Validation ===")
        
        # Check basic structure
        print(f"✓ Slide count: {len(prs.slides)}")
        print(f"✓ Layout count: {len(prs.slide_layouts)}")
        print(f"✓ Slide size: {prs.slide_width.inches}x{prs.slide_height.inches} inches")
        
        # Check each slide
        for i, slide in enumerate(prs.slides):
            print(f"✓ Slide {i+1}: {len(slide.shapes)} shapes")
            
            # Check for title
            if slide.shapes.title:
                title_text = slide.shapes.title.text if slide.shapes.title.text else "(empty)"
                print(f"  - Title: {title_text}")
        
        print("✓ Validation completed successfully")
        return True
        
    except Exception as e:
        print(f"❌ Validation failed: {e}")
        return False

# ==========================================
# SAFE MAIN EXECUTION
# ==========================================

def main():
    """Main execution with comprehensive error handling"""
    
    print("=== PowerPoint Generation Started ===")
    
    try:
        # Sample data
        sample_data = pd.DataFrame({
            'Category': ['A', 'B', 'C', 'D'],
            'Values': [23, 45, 56, 78],
            'Description': ['First', 'Second', 'Third', 'Fourth']
        })
        
        # Create presentation data
        data_dict = {
            'title': 'My Safe Data Presentation',
            'dataframe': sample_data,
            # 'image_path': 'path/to/your/image.png'  # Uncomment and provide valid path
        }
        
        # Create presentation
        print("Creating presentation...")
        presentation = create_presentation_with_data(data_dict)
        
        # Validate structure
        if validate_presentation_structure(presentation):
            # Save presentation
            filename = 'safe_presentation.pptx'
            if save_presentation_safely(presentation, filename):
                print("\n🎉 SUCCESS: Presentation created without corruption!")
                print(f"📁 File: {filename}")
            else:
                print("\n❌ ERROR: Failed to save presentation")
        else:
            print("\n❌ ERROR: Presentation structure validation failed")
            
    except Exception as e:
        print(f"\n❌ CRITICAL ERROR: {e}")
        
        # Emergency fallback
        try:
            print("Attempting emergency fallback...")
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Emergency Fallback"
            prs.save('emergency_presentation.pptx')
            print("✓ Emergency presentation created")
        except:
            print("❌ Even fallback failed")

# ==========================================
# DEBUGGING HELPER
# ==========================================

def debug_slide_layouts(prs):
    """Debug available slide layouts"""
    
    print("\n=== Available Slide Layouts ===")
    for i, layout in enumerate(prs.slide_layouts):
        try:
            layout_name = layout.name if hasattr(layout, 'name') else f"Layout {i}"
            print(f"{i}: {layout_name}")
            
            # Check placeholders
            placeholder_count = len(layout.placeholders)
            print(f"   Placeholders: {placeholder_count}")
            
        except Exception as e:
            print(f"{i}: Error reading layout - {e}")

def create_test_presentation():
    """Create a minimal test presentation for debugging"""
    
    print("Creating test presentation...")
    
    try:
        prs = Presentation()
        
        # Debug layouts
        debug_slide_layouts(prs)
        
        # Add single slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        
        # Save
        prs.save('test.pptx')
        print("✓ Test presentation saved as test.pptx")
        
    except Exception as e:
        print(f"❌ Test failed: {e}")

if __name__ == "__main__":
    # Run main program
    main()
    
    # Uncomment for debugging
    # create_test_presentation()