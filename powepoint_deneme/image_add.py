import os
from pptx import Presentation
from pptx.util import Inches

def test_image_addition():
    """Test adding an image to a PowerPoint slide"""
    
    try:
        print("Creating basic presentation...")
        
        # Create presentation
        prs = Presentation()
        
        # Add a slide (using title and content layout)
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        slide.shapes.title.text = "Image Test Slide"
        
        # Test image path - you can change this to any JPG image you have
        local_image_path = "test_image.jpg"  # Change this to your actual image path
        
        # Check if image exists
        if not os.path.exists(local_image_path):
            print(f"❌ Image file not found: {local_image_path}")
            print("Please provide a valid JPG image path or create a test_image.jpg file")
            return False
        
        # Image position and size settings
        left = Inches(2)        # Position from left
        top = Inches(2)         # Position from top  
        max_width = Inches(6)   # Maximum width
        
        print(f"Adding image: {local_image_path}")
        
        # YOUR CODE SNIPPET - Test this specific functionality
        pic = slide.shapes.add_picture(
            local_image_path, 
            left, 
            top, 
            width=max_width
        )
        
        print("✓ Image added successfully!")
        print(f"  - Image width: {pic.width.inches:.2f} inches")
        print(f"  - Image height: {pic.height.inches:.2f} inches")
        
        # Save the test presentation
        output_file = "test_image_presentation.pptx"
        prs.save(output_file)
        
        print(f"✓ Test presentation saved: {output_file}")
        
        # Check file size
        if os.path.exists(output_file):
            file_size = os.path.getsize(output_file)
            print(f"✓ File size: {file_size:,} bytes")
        
        return True
        
    except Exception as e:
        print(f"❌ Error: {e}")
        return False

if __name__ == "__main__":
    print("=== PowerPoint Image Addition Test ===")

    
    # Run the test
    success = test_image_addition()
    
    if success:
        print("\n🎉 SUCCESS! Image addition works correctly.")
        print("You can open 'test_image_presentation.pptx' to verify the result.")
    else:
        print("\n❌ FAILED! There's an issue with image addition.")
        
    print("\n=== Test Complete ===")